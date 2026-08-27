"""
FILE: routers/rag.py
─────────────────────────────────────────────────────────────────────────────
RAG Document-to-Quiz Router
MoSPI Skill Intelligence Platform | SIH 2026

Handles document ingestion (PDF, PPT, PPTX, TXT), validates payloads,
extracts cleaned plain-text using PyPDF / pdfplumber / python-pptx,
generates exactly 5 MCQs using Google Gemini API with strict JSON schema,
and provides quiz grading with automatic competency syncing to iGOT.
─────────────────────────────────────────────────────────────────────────────
"""

import io
import os
import re
import json
import uuid
import asyncio
import logging
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime, timezone
from typing import Optional, Dict, Any, Tuple, List

import httpx
from dotenv import load_dotenv
import google.generativeai as genai
from fastapi import APIRouter, File, HTTPException, UploadFile, status
from pydantic import BaseModel, Field

# Load environment variables (such as GEMINI_API_KEY, IGOT_COMPETENCIES_UPDATE_URL)
load_dotenv()

# PDF Extraction libraries
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import pypdf
except ImportError:
    pypdf = None

# PPTX Extraction library
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

router = APIRouter()

# ── Configuration Constants ───────────────────────────────────────────────────
MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024  # 25 MB max limit
SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".txt"}


# =============================================================================
# PYDANTIC SCHEMAS
# =============================================================================

class QuizQuestion(BaseModel):
    question: str = Field(..., description="The quiz question text")
    options: List[str] = Field(..., description="List of 4 multiple-choice options")
    correct_answer: int = Field(..., description="Zero-based index (0, 1, 2, or 3) of the correct option in options")
    explanation: str = Field(..., description="Detailed explanation for why the answer is correct")


class QuizPayload(BaseModel):
    questions: List[QuizQuestion] = Field(..., description="List of exactly 5 generated multiple-choice questions")


class DocumentMetadata(BaseModel):
    filename: str = Field(..., description="Original name of the uploaded file")
    file_type: str = Field(..., description="Detected file extension (.pdf, .ppt, .pptx, .txt)")
    file_size_bytes: int = Field(..., description="Size of the uploaded file in bytes")
    character_count: int = Field(..., description="Total extracted character count")
    word_count: int = Field(..., description="Total extracted word count")
    page_count: Optional[int] = Field(None, description="Total number of pages (for PDF)")
    slide_count: Optional[int] = Field(None, description="Total number of slides (for PPT/PPTX)")
    line_count: Optional[int] = Field(None, description="Total number of lines (for TXT)")


class DocumentUploadResponse(BaseModel):
    status: str = Field("success", description="Status code or status description ('success' or 'error')")
    message: str = Field(..., description="Human-readable status summary")
    quiz_id: str = Field(..., description="Unique identifier for the generated quiz session")
    filename: str = Field(..., description="Uploaded document filename")
    file_type: str = Field(..., description="Document file type")
    questions: List[QuizQuestion] = Field(..., description="List of 5 generated multiple choice questions based on document text")
    metadata: DocumentMetadata = Field(..., description="Structured metadata of the processed document")


class GradeRequest(BaseModel):
    user_id: str = Field(..., description="User ID or Government ID of the learner submitting the quiz")
    quiz_id: str = Field(..., description="Unique ID of the quiz session being graded")
    answers: List[int] = Field(..., description="List of chosen option indices corresponding to each question (0-indexed)")


class GradeResponse(BaseModel):
    status: str = Field("success", description="Status code or status description ('success' or 'error')")
    user_id: str = Field(..., description="User ID of the learner")
    quiz_id: str = Field(..., description="Quiz ID evaluated")
    score: float = Field(..., description="Percentage score achieved (0-100)")
    passed: bool = Field(..., description="True if score >= 70%, False otherwise")
    correct_count: int = Field(..., description="Number of correctly answered questions")
    total_questions: int = Field(..., description="Total number of questions in the quiz")
    message: str = Field(..., description="Human-readable result summary message")
    synced_to_igot: Optional[bool] = Field(None, description="Indicates whether the competency update was pushed to iGOT")
    igot_response: Optional[Dict[str, Any]] = Field(None, description="Response from the mock iGOT server if synced")


class ErrorResponse(BaseModel):
    status: str = Field("error", description="Failure indicator")
    detail: str = Field(..., description="Error message details")


# =============================================================================
# IN-MEMORY QUIZ STORE
# =============================================================================

QUIZ_STORE: Dict[str, Dict[str, Any]] = {
    "quiz_demo": {
        "quiz_id": "quiz_demo",
        "questions": [
            QuizQuestion(
                question="What is the primary international standard for National Accounts?",
                options=["SNA 2008", "GDP 1993", "IMF 2020", "OECD 2015"],
                correct_answer=0,
                explanation="System of National Accounts (SNA 2008) is the international statistical standard."
            ),
            QuizQuestion(
                question="Which organisation compiles CPI-Combined in India?",
                options=["Reserve Bank of India", "NSO / MoSPI", "NITI Aayog", "Ministry of Finance"],
                correct_answer=1,
                explanation="NSO (National Statistical Office) under MoSPI compiles and releases CPI."
            ),
            QuizQuestion(
                question="What is the base year for the current CPI series in India?",
                options=["2004-05", "2011-12", "2012", "2016"],
                correct_answer=2,
                explanation="The current CPI series has base year 2012 = 100."
            ),
            QuizQuestion(
                question="Which sampling technique gives every unit an equal probability of selection?",
                options=["Simple Random Sampling", "Judgment Sampling", "Quota Sampling", "Snowball Sampling"],
                correct_answer=0,
                explanation="Simple Random Sampling (SRS) ensures equal inclusion probability."
            ),
            QuizQuestion(
                question="What formula is predominantly used for consumer price indices in India?",
                options=["Laspeyres formula", "Fisher ideal index", "Törnqvist index", "Divisia index"],
                correct_answer=0,
                explanation="Laspeyres base-weighted formula is standard in CPI compilation."
            )
        ],
        "filename": "sample_mospi_overview.pdf",
        "competency_id": "FRAC-STAT-001",
        "skill_name": "National Accounts & Official Statistics",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
}


# =============================================================================
# TEXT EXTRACTION HELPERS
# =============================================================================

def _clean_text(text: str) -> str:
    """
    Normalizes whitespace, removes null bytes, and cleans up repeated linebreaks.
    """
    if not text:
        return ""
    # Remove null bytes
    text = text.replace("\x00", "")
    # Normalize carriage returns
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Replace multiple continuous whitespace (except single newlines)
    text = re.sub(r"[ \t]+", " ", text)
    # Condense 3+ newlines to double newline
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pdf(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text from PDF bytes using pdfplumber (primary) with fallback to pypdf.
    Returns a tuple of (extracted_text, page_count).
    """
    extracted_pages = []
    page_count = 0

    # 1. Attempt extraction with pdfplumber (superior layout & table preservation)
    if pdfplumber is not None:
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                page_count = len(pdf.pages)
                for idx, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        extracted_pages.append(f"--- Page {idx} ---\n{page_text.strip()}")
            
            combined_text = "\n\n".join(extracted_pages)
            if combined_text.strip():
                return _clean_text(combined_text), page_count
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}. Falling back to pypdf.")
            extracted_pages.clear()

    # 2. Fallback to pypdf
    if pypdf is not None:
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            page_count = len(reader.pages)
            for idx, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    extracted_pages.append(f"--- Page {idx} ---\n{page_text.strip()}")
            
            combined_text = "\n\n".join(extracted_pages)
            return _clean_text(combined_text), page_count
        except Exception as e:
            logger.error(f"pypdf extraction failed: {e}")
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Could not parse PDF content: {str(e)}"
            )

    raise HTTPException(
        status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
        detail="No PDF extraction library available on the server (pdfplumber/pypdf)."
    )


def _extract_pptx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text from PPTX files using python-pptx with XML fallback.
    Returns a tuple of (extracted_text, slide_count).
    """
    slide_texts = []
    slide_count = 0

    # 1. Primary extraction via python-pptx
    if Presentation is not None:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            slide_count = len(prs.slides)

            for idx, slide in enumerate(prs.slides, start=1):
                parts = []
                # Extract slide shape texts (titles, textboxes, tables)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in paragraph.runs).strip()
                            if text:
                                parts.append(text)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                parts.append(row_text)

                # Extract speaker notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes: {notes}]")

                if parts:
                    slide_texts.append(f"--- Slide {idx} ---\n" + "\n".join(parts))

            combined_text = "\n\n".join(slide_texts)
            if combined_text.strip():
                return _clean_text(combined_text), slide_count
        except Exception as e:
            logger.warning(f"python-pptx failed: {e}. Attempting direct XML parsing.")
            slide_texts.clear()

    # 2. Fallback: Parse slide XML files from PPTX (ZIP container) directly
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            slide_files = [f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")]
            slide_count = len(slide_files)
            # Sort numerically by slide number
            slide_files.sort(key=lambda name: int(re.search(r"\d+", name).group()) if re.search(r"\d+", name) else 0)

            for idx, slide_file in enumerate(slide_files, start=1):
                xml_content = z.read(slide_file)
                tree = ET.fromstring(xml_content)
                # PPTX text elements are in the 'a:t' tags
                texts = [elem.text for elem in tree.iter() if elem.tag.endswith("}t") and elem.text]
                if texts:
                    slide_texts.append(f"--- Slide {idx} ---\n" + "\n".join(texts))

            combined_text = "\n\n".join(slide_texts)
            return _clean_text(combined_text), slide_count
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Could not parse PowerPoint presentation: {str(e)}"
        )


def _extract_txt(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text from plain text bytes trying multiple character encodings.
    Returns a tuple of (extracted_text, line_count).
    """
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]
    decoded_text = None

    for enc in encodings:
        try:
            decoded_text = file_bytes.decode(enc)
            break
        except UnicodeDecodeError:
            continue

    if decoded_text is None:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="Unable to decode text file. Ensure it is encoded in UTF-8 or standard ASCII."
        )

    cleaned = _clean_text(decoded_text)
    line_count = len(decoded_text.splitlines())
    return cleaned, line_count


# =============================================================================
# GEMINI QUIZ GENERATION HELPER
# =============================================================================

async def _generate_mcqs_from_text(text: str) -> List[QuizQuestion]:
    """
    Calls Google Gemini API using google-generativeai with GEMINI_API_KEY from .env.
    Generates exactly 5 MCQs based on the extracted text and enforces strict JSON output schema:
    {"questions": [{"question": "...", "options": ["A", "B", "C", "D"], "correct_answer": 0, "explanation": "..."}]}
    """
    load_dotenv()
    api_key = os.getenv("GEMINI_API_KEY", "").strip()

    if not api_key or api_key in {"your-actual-api-key-here", "YOUR_GEMINI_API_KEY"}:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="GEMINI_API_KEY is not configured in the .env file. Please configure a valid Gemini API key.",
        )

    try:
        genai.configure(api_key=api_key)
        
        configured_model = os.getenv("GEMINI_MODEL", "gemini-3.5-flash-lite").strip()
        candidate_models = [configured_model] if configured_model else []
        for fallback in ["gemini-3.5-flash-lite", "gemini-3.6-flash", "gemini-3.5-flash", "gemini-flash-latest"]:
            if fallback not in candidate_models:
                candidate_models.append(fallback)

        prompt = (
            "You are an expert assessment and quiz creation engine for the Ministry of Statistics and Programme Implementation (MoSPI).\n"
            "Generate exactly 5 multiple-choice questions (MCQs) strictly based on the following training document content.\n\n"
            "Strict Requirements:\n"
            "1. Generate exactly 5 high-quality questions testing comprehension, core concepts, or statistical procedures described in the text.\n"
            "2. Each question MUST contain exactly 4 distinct and plausible options.\n"
            "3. 'correct_answer' MUST be the 0-based integer index of the correct option in the 'options' array (0, 1, 2, or 3).\n"
            "4. Provide a clear, detailed 'explanation' justifying why the selected option is correct according to the text.\n"
            "5. Return output conforming strictly to the JSON schema: {\"questions\": [{\"question\": \"...\", \"options\": [\"A\", \"B\", \"C\", \"D\"], \"correct_answer\": 0, \"explanation\": \"...\"}]}.\n\n"
            f"Document Content:\n\"\"\"\n{text[:40000]}\n\"\"\""
        )

        response = None
        last_error = None

        for model_name in candidate_models:
            try:
                model = genai.GenerativeModel(
                    model_name=model_name,
                    generation_config={
                        "response_mime_type": "application/json",
                        "temperature": 0.3,
                    },
                )
                response = await asyncio.to_thread(model.generate_content, prompt)
                if response and response.text:
                    break
            except Exception as ex:
                logger.warning(f"Gemini model '{model_name}' attempt failed: {ex}")
                last_error = ex
                continue

        if not response or not response.text:
            raise ValueError(f"All Gemini model attempts failed. Last error: {last_error}")

        raw_text = response.text.strip()

        # Clean code fence blocks if returned by model
        cleaned_text = raw_text
        if cleaned_text.startswith("```json"):
            cleaned_text = cleaned_text[7:]
        elif cleaned_text.startswith("```"):
            cleaned_text = cleaned_text[3:]
        if cleaned_text.endswith("```"):
            cleaned_text = cleaned_text[:-3]
        cleaned_text = cleaned_text.strip()

        parsed_json = json.loads(cleaned_text)
        if isinstance(parsed_json, dict) and "questions" in parsed_json:
            raw_questions = parsed_json["questions"]
        elif isinstance(parsed_json, list):
            raw_questions = parsed_json
        else:
            raw_questions = [parsed_json]

        normalized_questions: List[QuizQuestion] = []
        for item in raw_questions:
            if not isinstance(item, dict):
                continue
            
            q_text = str(item.get("question") or item.get("question_text") or item.get("prompt") or "").strip()
            if not q_text:
                continue

            raw_opts = item.get("options") or item.get("choices") or item.get("answers") or []
            if isinstance(raw_opts, dict):
                opts = [str(v).strip() for v in raw_opts.values()]
            elif isinstance(raw_opts, list):
                opts = [str(o).strip() for o in raw_opts]
            else:
                opts = [str(raw_opts).strip()]

            opts = [o for o in opts if o]
            while len(opts) < 4:
                opts.append(f"Option {chr(65 + len(opts))}")
            if len(opts) > 4:
                opts = opts[:4]

            raw_ans = item.get("correct_answer") if "correct_answer" in item else (item.get("correctAnswer") if "correctAnswer" in item else item.get("answer"))
            ans_idx = 0
            if isinstance(raw_ans, int):
                ans_idx = max(0, min(len(opts) - 1, raw_ans))
            elif isinstance(raw_ans, str):
                ans_str = raw_ans.strip().upper()
                if ans_str in {"A", "B", "C", "D"}:
                    ans_idx = ord(ans_str) - ord("A")
                elif ans_str.isdigit():
                    ans_idx = max(0, min(len(opts) - 1, int(ans_str)))
                else:
                    for i, opt in enumerate(opts):
                        if opt.lower() == raw_ans.strip().lower():
                            ans_idx = i
                            break

            explanation = str(
                item.get("explanation")
                or item.get("reason")
                or item.get("justification")
                or f"Option {chr(65 + ans_idx)} is the correct answer according to the document."
            ).strip()

            normalized_questions.append(
                QuizQuestion(
                    question=q_text,
                    options=opts,
                    correct_answer=ans_idx,
                    explanation=explanation
                )
            )

        if not normalized_questions:
            raise ValueError("No valid questions could be extracted from Gemini response.")

        return normalized_questions

    except HTTPException:
        raise
    except json.JSONDecodeError as jde:
        logger.error(f"JSON decode failure from Gemini response: {jde}")
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Failed to parse structured JSON quiz questions from Gemini response: {str(jde)}",
        )
    except Exception as e:
        logger.exception(f"Gemini API generation error: {e}")
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Gemini API error during quiz generation: {str(e)}",
        )


# =============================================================================
# ENDPOINTS
# =============================================================================

@router.post(
    "/upload",
    response_model=DocumentUploadResponse,
    status_code=status.HTTP_200_OK,
    summary="Upload document, extract text, and generate 5 RAG Quiz MCQs with Gemini API",
    responses={
        200: {"description": "Document parsed and 5 MCQs successfully generated via Gemini API."},
        400: {"model": ErrorResponse, "description": "Invalid file format or empty payload."},
        413: {"model": ErrorResponse, "description": "File exceeds maximum size limit."},
        422: {"model": ErrorResponse, "description": "Unprocessable document or no readable text."},
        500: {"model": ErrorResponse, "description": "Internal server processing or configuration error."},
        502: {"model": ErrorResponse, "description": "Gemini API error or bad gateway during generation."},
    },
)
async def upload_document_for_rag(
    file: UploadFile = File(..., description="Document file to process (.pdf, .ppt, .pptx, .txt)")
) -> DocumentUploadResponse:
    """
    **RAG Document Ingestion & Quiz Generator**
    
    - Accepts **PDF**, **PPT/PPTX**, and **TXT** files.
    - Validates file extension, MIME type, and size constraints (max 25MB).
    - Extracts clean plain text with page/slide context using **pdfplumber / PyPDF / python-pptx**.
    - Calls Google Gemini API using `google-generativeai` with `GEMINI_API_KEY` from `.env`.
    - Generates exactly **5 MCQs** with strict JSON schema:
      `{"questions": [{"question": "...", "options": ["A", "B", "C", "D"], "correct_answer": 0, "explanation": "..."}]}`.
    - Stores generated quiz in-memory with unique `quiz_id` for grading.
    - Returns the structured quiz payload, quiz_id, and document metadata.
    """
    # 1. Validate presence of filename
    if not file.filename:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No filename provided in upload payload."
        )

    filename = os.path.basename(file.filename)
    _, ext = os.path.splitext(filename.lower())

    # 2. Validate file extension
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=(
                f"Unsupported file format '{ext}'. "
                f"Allowed document formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )
        )

    # 3. Read file content safely
    try:
        file_bytes = await file.read()
    except Exception as e:
        logger.error(f"Error reading uploaded stream for {filename}: {e}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to read uploaded file: {str(e)}"
        )
    finally:
        await file.close()

    # 4. Validate file size
    file_size = len(file_bytes)
    if file_size == 0:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="The uploaded file is empty (0 bytes)."
        )

    if file_size > MAX_FILE_SIZE_BYTES:
        max_mb = MAX_FILE_SIZE_BYTES // (1024 * 1024)
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail=f"File size exceeds maximum allowed limit of {max_mb} MB."
        )

    # 5. Extract text based on file type
    page_count: Optional[int] = None
    slide_count: Optional[int] = None
    line_count: Optional[int] = None
    extracted_text = ""

    try:
        if ext == ".pdf":
            extracted_text, page_count = _extract_pdf(file_bytes)
        elif ext in {".ppt", ".pptx"}:
            if ext == ".ppt":
                # Check if it's actually a PPTX named .ppt or legacy PPT
                if file_bytes.startswith(b"PK\x03\x04"):
                    extracted_text, slide_count = _extract_pptx(file_bytes)
                else:
                    raise HTTPException(
                        status_code=status.HTTP_400_BAD_REQUEST,
                        detail="Legacy binary .ppt files must be converted to .pptx format or PDF."
                    )
            else:
                extracted_text, slide_count = _extract_pptx(file_bytes)
        elif ext == ".txt":
            extracted_text, line_count = _extract_txt(file_bytes)

    except HTTPException:
        raise
    except Exception as e:
        logger.exception(f"Unexpected error extracting text from {filename}: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Internal error processing document: {str(e)}"
        )

    # 6. Verify that readable text was obtained
    if not extracted_text or len(extracted_text.strip()) == 0:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=(
                "No readable text could be extracted from this document. "
                "If this is a scanned PDF/image, please provide a document with OCR/selectable text."
            )
        )

    # 7. Generate exactly 5 MCQs via Google Gemini API
    questions = await _generate_mcqs_from_text(extracted_text)

    # 8. Generate unique quiz ID and store in-memory for subsequent grading
    quiz_id = f"QZ-{uuid.uuid4().hex[:8].upper()}"
    QUIZ_STORE[quiz_id] = {
        "quiz_id": quiz_id,
        "questions": questions,
        "filename": filename,
        "competency_id": "FRAC-STAT-001",
        "skill_name": "National Accounts & Official Statistics",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }

    # 9. Compute statistics & build response metadata
    char_count = len(extracted_text)
    word_count = len(extracted_text.split())

    metadata = DocumentMetadata(
        filename=filename,
        file_type=ext,
        file_size_bytes=file_size,
        character_count=char_count,
        word_count=word_count,
        page_count=page_count,
        slide_count=slide_count,
        line_count=line_count,
    )

    return DocumentUploadResponse(
        status="success",
        message=f"Successfully generated {len(questions)} quiz questions from {filename}.",
        quiz_id=quiz_id,
        filename=filename,
        file_type=ext,
        questions=questions,
        metadata=metadata,
    )


@router.post(
    "/grade",
    response_model=GradeResponse,
    status_code=status.HTTP_200_OK,
    summary="Grade submitted quiz answers and sync competency update to iGOT if passed",
    responses={
        200: {"description": "Quiz evaluated successfully."},
        400: {"model": ErrorResponse, "description": "Invalid submission or answers format."},
        404: {"model": ErrorResponse, "description": "Quiz session not found."},
        500: {"model": ErrorResponse, "description": "Internal grading error."},
    },
)
async def grade_quiz(payload: GradeRequest) -> GradeResponse:
    """
    **Grade Quiz & Sync Competency to Mock iGOT Server**
    
    1. Accepts JSON payload with `user_id`, `quiz_id`, and `answers` (list of selected option indices).
    2. Compares submitted option indices against stored correct answers.
    3. Calculates percentage score: `(correct_count / total_questions) * 100`.
    4. If `score >= 70%` (Pass):
       - Calls external mock iGOT server at `http://localhost:8001/competencies/update` using `httpx`.
       - Transmits updated skill / competency payload.
    5. Returns score, pass/fail status, detailed message, and iGOT synchronization results.
    """
    # 1. Validate quiz existence in store
    quiz = QUIZ_STORE.get(payload.quiz_id)
    if not quiz:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=f"Quiz session '{payload.quiz_id}' not found. Please upload a document to generate a quiz or provide a valid quiz_id.",
        )

    questions: List[QuizQuestion] = quiz.get("questions", [])
    total_questions = len(questions)

    if total_questions == 0:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="The specified quiz contains no questions to grade.",
        )

    # 2. Validate answers list
    if payload.answers is None or len(payload.answers) == 0:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No answers provided in submission payload.",
        )

    # 3. Calculate score
    correct_count = 0
    for idx, question in enumerate(questions):
        if idx < len(payload.answers):
            user_choice = payload.answers[idx]
            if user_choice == question.correct_answer:
                correct_count += 1

    score_percentage = round((correct_count / total_questions) * 100.0, 2)
    passed = score_percentage >= 70.0

    # 4. Sync with mock iGOT server if passed (score >= 70%)
    synced_to_igot: Optional[bool] = None
    igot_response_data: Optional[Dict[str, Any]] = None

    if passed:
        igot_url = os.getenv("IGOT_COMPETENCIES_UPDATE_URL", "http://localhost:8001/competencies/update")
        igot_token = os.getenv("IGOT_MOCK_TOKEN", "mock-api-key-2026")

        update_payload = {
            "user_id": payload.user_id,
            "competency": "Data Analysis",
            "new_level": 3,
        }

        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                resp = await client.post(
                    igot_url,
                    json=update_payload,
                    headers={"x-authenticated-user-token": igot_token},
                )
                if resp.status_code in (200, 201):
                    synced_to_igot = True
                    try:
                        igot_response_data = resp.json()
                    except Exception:
                        igot_response_data = {"status": "success", "raw": resp.text}
                else:
                    synced_to_igot = False
                    logger.warning(f"iGOT competency update responded with HTTP {resp.status_code}: {resp.text}")
                    igot_response_data = {"status_code": resp.status_code, "response": resp.text}

        except httpx.RequestError as exc:
            synced_to_igot = False
            logger.warning(f"Could not connect to mock iGOT server at {igot_url}: {exc}")
            igot_response_data = {"warning": f"iGOT server unreachable at {igot_url}: {str(exc)}"}
        except Exception as exc:
            synced_to_igot = False
            logger.error(f"Unexpected error communicating with mock iGOT server: {exc}")
            igot_response_data = {"error": str(exc)}

    # 5. Formulate summary message
    if passed:
        msg = f"Passed! You scored {score_percentage}% ({correct_count}/{total_questions} correct)."
        if synced_to_igot:
            msg += " Competency record was successfully updated on the iGOT platform."
        elif synced_to_igot is False:
            msg += " (Note: iGOT mock server sync was not reachable, but assessment was passed)."
    else:
        msg = f"Did not pass. You scored {score_percentage}% ({correct_count}/{total_questions} correct). A minimum score of 70% is required to pass and update your competency profile."

    return GradeResponse(
        status="success",
        user_id=payload.user_id,
        quiz_id=payload.quiz_id,
        score=score_percentage,
        passed=passed,
        correct_count=correct_count,
        total_questions=total_questions,
        message=msg,
        synced_to_igot=synced_to_igot,
        igot_response=igot_response_data,
    )
